import os
import sys

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

file_path = r"d:\demand-planning\reference data\Demand_Planning_Engine_BRD.pptx"

try:
    prs = Presentation(file_path)
    with open("brd_text.txt", "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
            f.write("\n\n")
    print("Successfully extracted to brd_text.txt")
except Exception as e:
    print(f"Error reading PPTX: {e}")
